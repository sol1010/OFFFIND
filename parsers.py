"""엑셀(xlsx/xlsm) / PDF / Word(docx) / PowerPoint(pptx) 파일에서 검색 가능한
텍스트 조각(entry)을 추출. 구버전 바이너리 포맷(.xls/.doc/.ppt)은 대상 라이브러리가
아예 못 읽어서 지원하지 않는다(README 알려진 제약에 명시)."""
from typing import List, Dict, Any


class ParseFailed(Exception):
    """파일을 아예 열지 못했다(손상·잠김·지원 안 되는 변형 등).

    "열긴 했는데 글자가 없다"(예: 스캔 이미지만 있는 PDF, 빈 문서)와 반드시
    구분해야 해서 예외로 올린다. 예전엔 둘 다 빈 리스트를 돌려줘서 호출자
    (indexer)가 구분할 수 없었고, 그래서 색인 순간에 잠깐 열려 있던 파일이
    "내용 없는 파일"로 굳어 영원히 검색에서 빠졌다.

    반대로 파일은 열렸는데 시트/페이지 하나에서만 실패한 경우는 예외로 올리지
    않는다 — 그때까지 뽑은 내용은 쓸모가 있고, 다시 시도해도 어차피 같은
    지점에서 또 실패하기 때문이다."""


def extract_xlsx(path: str) -> List[Dict[str, Any]]:
    """시트별 각 행을 하나의 검색 entry로 반환."""
    import openpyxl

    entries = []
    try:
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    except Exception as e:
        raise ParseFailed(f"xlsx 열기 실패: {path}") from e

    try:
        # docx/pptx 추출기와 마찬가지로, 시트/행 하나에서 예상 못한 예외가 나도
        # 그때까지 뽑은 entries는 살린다 — 이 감싸기가 없으면 예외가 finally만
        # 거치고 return entries 줄 자체를 건너뛰어서, 앞선 시트들에서 이미 잘
        # 뽑아둔 내용까지 호출자(rebuild())의 바깥쪽 except가 통째로 버리게 된다
        # (pptx 노트에서 실제로 겪었던 것과 같은 종류의 항목별 파싱 실패).
        try:
            for sheet in wb.worksheets:
                for row_idx, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                    cells = [str(c) for c in row if c is not None and str(c).strip() != ""]
                    if not cells:
                        continue
                    text = " | ".join(cells)
                    entries.append({
                        "location": f"{sheet.title} · {row_idx}행",
                        "text": text,
                        "sheet": sheet.title,
                        "row": row_idx,
                    })
        except Exception:
            pass
    finally:
        wb.close()
    return entries


def extract_pdf(path: str) -> List[Dict[str, Any]]:
    """페이지별 텍스트를 하나의 검색 entry로 반환."""
    import pdfplumber

    entries = []
    try:
        pdf = pdfplumber.open(path)
    except Exception as e:
        raise ParseFailed(f"pdf 열기 실패: {path}") from e

    # 열린 뒤 페이지 하나에서 나는 실패는 그때까지 뽑은 것을 살린다(ParseFailed 참고).
    try:
        try:
            for i, page in enumerate(pdf.pages, start=1):
                text = page.extract_text() or ""
                text = " ".join(text.split())
                if not text:
                    continue
                entries.append({
                    "location": f"{i}페이지",
                    "text": text,
                    "page": i,
                })
        except Exception:
            pass
    finally:
        pdf.close()
    return entries


def extract_docx(path: str) -> List[Dict[str, Any]]:
    """문단 단위로 텍스트를 추출한다(표 안 텍스트는 표 단위로 별도 entry)."""
    import docx

    entries = []
    try:
        doc = docx.Document(path)
    except Exception as e:
        raise ParseFailed(f"docx 열기 실패: {path}") from e

    # 개별 문단/표에서 예상 못한 예외(문서마다 구조가 미묘하게 다를 수 있다)가
    # 나도 전체 파일을 통째로 포기하지 않고, 그때까지 뽑은 entries 는 그대로
    # 쓴다 — extract_pdf 가 페이지 단위 실패에 관대한 것과 같은 이유다. rebuild()
    # 가 백그라운드 스레드에서 도는데, 여기서 예외가 새 나가면 그 스레드 전체가
    # 죽어서 색인이 중간에 멈춰버린다(실제로 한 번 이렇게 겪음).
    try:
        for i, para in enumerate(doc.paragraphs, start=1):
            text = para.text.strip()
            if not text:
                continue
            entries.append({
                "location": f"{i}번째 문단",
                "text": text,
                "paragraph": i,
            })

        # 표 안 텍스트는 doc.paragraphs 에 안 잡힌다(별도 요소라서) — 놓치면 안 되니
        # 표는 표대로 행 단위로 추출한다. 표 안에서 더블클릭하면 문서는 열리지만
        # 문단 이동과 달리 그 행으로 바로 스크롤하는 기능은 없다(paragraph 키가 없음).
        for t_idx, table in enumerate(doc.tables, start=1):
            for r_idx, row in enumerate(table.rows, start=1):
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if not cells:
                    continue
                entries.append({
                    "location": f"표{t_idx} · {r_idx}행",
                    "text": " | ".join(cells),
                })
    except Exception:
        pass
    return entries


def extract_pptx(path: str) -> List[Dict[str, Any]]:
    """슬라이드 단위로 텍스트를 추출한다(표/발표자 노트 포함)."""
    import pptx

    entries = []
    try:
        prs = pptx.Presentation(path)
    except Exception as e:
        raise ParseFailed(f"pptx 열기 실패: {path}") from e

    # 슬라이드 하나에서 예상 못한 예외가 나도 전체 파일을 포기하지 않는다 —
    # 실제로 "노트 슬라이드는 있는데(has_notes_slide=True) 그 안에 노트
    # placeholder 도형 자체가 없어서 notes_text_frame 이 None"인 경우가 있었고
    # (AttributeError), 이게 rebuild() 를 도는 백그라운드 스레드를 통째로
    # 죽여서 색인이 중간에 멈췄다.
    try:
        for i, slide in enumerate(prs.slides, start=1):
            parts = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    parts.append(shape.text_frame.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            parts.append(" | ".join(cells))
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(notes)
            text = " ".join(" ".join(parts).split())
            if not text:
                continue
            entries.append({
                "location": f"{i}번 슬라이드",
                "text": text,
                "slide": i,
            })
    except Exception:
        pass
    return entries


EXTRACTORS = {
    ".xlsx": extract_xlsx,
    ".xlsm": extract_xlsx,
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
}
